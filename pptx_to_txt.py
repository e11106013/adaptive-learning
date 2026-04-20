#華語教材 PPTX 轉 TXT 工具
#使用目的：將華語教材簡報檔案 .pptx 轉為文字檔並加上 YAML

#命名規則：B{冊數}-L{課數}-{類別}.pptx
#輸出格式：帶 YAML 標頭的 .txt 檔案

#使用方式：
# python pptx_to_txt.py --input C:\Users\ROG\Documents\PPT2txt\raw --output C:\Users\ROG\Documents\PPT2txt\txt

#安裝需求：
# pip install python-pptx pyyaml

import os
import re
import argparse
from pptx import Presentation
import yaml


def parse_filename(filename):
    """
    從檔名解析冊數、課數、類別
    格式：B3-L01-Grammar.pptx
    """
    pattern = r"B(\d+)-L(\d+)-(\w+)\.pptx$"
    match = re.match(pattern, filename, re.IGNORECASE)
    if not match:
        return None
    book = int(match.group(1))
    lesson = int(match.group(2))
    type_raw = match.group(3).lower()

    # 標準化類別名稱
    if "grammar" in type_raw:
        content_type = "grammar"
    elif "vocab" in type_raw:
        content_type = "vocabulary"
    else:
        content_type = type_raw

    return {
        "book": book,
        "lesson": lesson,
        "type": content_type,
        "level": f"B{book}-L{lesson:02d}"
    }


def extract_text_from_pptx(filepath):
    """
    從 PPTX 抽取所有文字
    依投影片順序，每張投影片之間用空行分隔
    """
    prs = Presentation(filepath)
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    slide_lines.append(line)
        if slide_lines:
            slides_text.append("\n".join(slide_lines))

    return "\n\n".join(slides_text)


def build_yaml_header(meta):
    """
    建立 YAML 標頭
    """
    header = {
        "book": meta["book"],
        "lesson": meta["lesson"],
        "level": meta["level"],
        "type": meta["type"]
    }
    return "---\n" + yaml.dump(header, allow_unicode=True, default_flow_style=False) + "---\n"


def convert_pptx_to_txt(input_dir, output_dir):
    """
    掃描資料夾，將所有符合命名規則的 PPTX 轉換為 TXT
    """
    os.makedirs(output_dir, exist_ok=True)

    pptx_files = [f for f in os.listdir(input_dir) if f.endswith(".pptx")]
    if not pptx_files:
        print("找不到任何 .pptx 檔案")
        return

    success = 0
    skipped = 0

    for filename in sorted(pptx_files):
        meta = parse_filename(filename)
        if not meta:
            print(f"[跳過] 檔名格式不符：{filename}")
            skipped += 1
            continue

        filepath = os.path.join(input_dir, filename)
        content = extract_text_from_pptx(filepath)

        if not content.strip():
            print(f"[警告] 沒有抽取到文字：{filename}")
            skipped += 1
            continue

        yaml_header = build_yaml_header(meta)
        output_content = yaml_header + "\n" + content

        output_filename = filename.replace(".pptx", ".txt")
        output_path = os.path.join(output_dir, output_filename)

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(output_content)

        print(f"[完成] {filename} → {output_filename}")
        success += 1

    print(f"\n完成：{success} 個檔案轉換成功，{skipped} 個跳過")


def main():
    parser = argparse.ArgumentParser(description="華語教材 PPTX 轉 TXT 工具")
    parser.add_argument("--input",  required=True, help="PPTX 檔案所在的資料夾")
    parser.add_argument("--output", required=True, help="TXT 輸出的資料夾")
    args = parser.parse_args()

    convert_pptx_to_txt(args.input, args.output)


if __name__ == "__main__":
    main()
